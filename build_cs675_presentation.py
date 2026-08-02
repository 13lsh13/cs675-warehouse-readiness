#!/usr/bin/env python3
"""Build the CS-675 final presentation from the generated sample data."""

import csv
from collections import Counter, defaultdict
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
DATA = ROOT / "data" / "sample"
ART = ROOT / "artifacts"
ART.mkdir(exist_ok=True)

NAVY = "102A43"
BLUE = "1976D2"
CYAN = "00A6A6"
ORANGE = "F39C3D"
RED = "D64550"
GREEN = "2E8B57"
INK = "243B53"
MUTED = "627D98"
LIGHT = "F4F8FB"
WHITE = "FFFFFF"


def read_csv(name):
    with (DATA / name).open(encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


products = read_csv("product_reference.csv")
shipments = read_csv("shipment_schedule.csv")
transactions = read_csv("inventory_transactions.csv")
readiness = read_csv("order_readiness_expected.csv")
statuses = Counter(row["readiness_status"] for row in readiness)
missing_total = sum(int(row["missing_pallets"]) for row in readiness)


def create_charts():
    plt.style.use("seaborn-v0_8-whitegrid")
    order = ["READY", "AT_RISK", "WAITING_FOR_INBOUND", "SHORT"]
    labels = ["Ready", "At Risk", "Waiting for\nInbound", "Short"]
    colors = [f"#{GREEN}", f"#{ORANGE}", f"#{BLUE}", f"#{RED}"]
    fig, ax = plt.subplots(figsize=(8.4, 4.1))
    values = [statuses[x] for x in order]
    bars = ax.bar(labels, values, color=colors, width=.65)
    ax.set_ylabel("Upcoming outbound orders")
    ax.set_title("Most orders are ready, but 173 require attention", loc="left", weight="bold", color=f"#{NAVY}")
    ax.spines[["top", "right"]].set_visible(False)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, value + 8, f"{value}", ha="center", weight="bold")
    fig.tight_layout()
    fig.savefig(ART / "readiness_status.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    volume = defaultdict(lambda: [0, 0])
    for row in shipments:
        if row["shipment_status"] == "CANCELLED":
            continue
        idx = 0 if row["shipment_type"] == "INBOUND" else 1
        volume[row["product_id"]][idx] += int(row["planned_pallet_quantity"])
    top = sorted(volume.items(), key=lambda x: sum(x[1]), reverse=True)[:10][::-1]
    fig, ax = plt.subplots(figsize=(8.4, 4.1))
    names = [x[0] for x in top]
    inbound = [x[1][0] for x in top]
    outbound = [x[1][1] for x in top]
    ax.barh(names, inbound, label="Inbound", color=f"#{CYAN}")
    ax.barh(names, outbound, left=inbound, label="Outbound", color=f"#{NAVY}")
    ax.set_xlabel("Planned pallets")
    ax.set_title("Top products by planned pallet movement", loc="left", weight="bold", color=f"#{NAVY}")
    ax.legend(frameon=False, ncol=2, loc="lower right")
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    fig.savefig(ART / "product_volume.png", dpi=180, bbox_inches="tight")
    plt.close(fig)


create_charts()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def background(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = rgb(color)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(.16), prs.slide_height)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = rgb(CYAN); stripe.line.fill.background()


def text(slide, value, x, y, w, h, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = value; p.alignment = align
    p.font.name = "Aptos"; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = rgb(color)
    return shape


def heading(slide, title, kicker):
    text(slide, kicker.upper(), .62, .25, 12, .3, 10, CYAN, True)
    text(slide, title, .62, .58, 12, .65, 27, NAVY, True)


def bullets(slide, items, x, y, w, h, size=18):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item; p.space_after = Pt(12)
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = rgb(INK)
    return shape


def card(slide, title_value, body, x, y, w, h, accent=BLUE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = rgb(WHITE); box.line.color.rgb = rgb("D9E2EC")
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(.08), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(accent); bar.line.fill.background()
    text(slide, title_value, x+.25, y+.13, w-.42, .4, 15, NAVY, True)
    text(slide, body, x+.25, y+.58, w-.42, h-.7, 13, MUTED)


def add_notes(slide, note):
    try:
        slide.notes_slide.notes_text_frame.text = note
    except Exception:
        pass


def footer(slide, number):
    text(slide, f"CS-675 Big Data Management  |  Shuhao Lin", .63, 7.16, 5.2, .2, 9, MUTED)
    text(slide, str(number), 12.25, 7.16, .4, .2, 9, MUTED, align=PP_ALIGN.RIGHT)


def new_slide(title_value, kicker, note):
    slide = prs.slides.add_slide(blank); background(slide); heading(slide, title_value, kicker)
    footer(slide, len(prs.slides)); add_notes(slide, note)
    return slide


# 1
s = prs.slides.add_slide(blank); background(s, NAVY)
text(s, "WAREHOUSE INVENTORY &\nOUTBOUND ORDER READINESS", .7, 1.2, 8.2, 1.75, 30, WHITE, True)
text(s, "A scalable Spark and AWS system for identifying pallet shortages before loading", .72, 3.18, 7.5, .8, 18, "B8D8EA")
card(s, "DECISION QUESTION", "Will enough inventory exist for each outbound order—and if not, when will it become available?", 8.9, 1.4, 3.65, 2.25, CYAN)
text(s, "Shuhao Lin  •  CS-675 Big Data Management", .72, 6.55, 7, .3, 13, WHITE)
add_notes(s, "Good afternoon. My final project is a Warehouse Inventory and Outbound Order Readiness Analysis System. It combines warehouse operations knowledge with Spark, SQL, AWS, and reproducible infrastructure.")

# 2
s = new_slide("Inventory on hand is only part of the answer", "Business problem", "Warehouses must consider timing, not only the current balance. An order can look short today but become ready after an inbound receipt—or look ready until an earlier outbound consumes the stock.")
card(s, "Current inventory", "Completed receipts, shipments, transfers, and adjustments establish the as-of balance.", .7, 1.55, 3.7, 1.6, BLUE)
card(s, "Future events", "Scheduled inbound and outbound quantities change the projected balance chronologically.", 4.8, 1.55, 3.7, 1.6, CYAN)
card(s, "Operational decision", "Flag the order, quantify missing pallets, and estimate the first available date.", 8.9, 1.55, 3.7, 1.6, ORANGE)
text(s, "PROJECT QUESTION", .75, 3.7, 2.3, .3, 11, CYAN, True)
text(s, "Will the warehouse have enough inventory to complete each upcoming outbound order?", .75, 4.05, 11.7, .95, 25, NAVY, True)
text(s, "If not: how many pallets are missing, and when is inventory expected to become available?", .75, 5.2, 11.5, .7, 18, INK)

# 3
s = new_slide("Three related datasets create an event-level inventory view", "Data model", "The design separates product master data, shipment plans, and executed inventory events. Product ID and warehouse ID connect the tables; shipment ID links scheduled shipments to completed movements.")
card(s, "PRODUCT REFERENCE", "product_id • category • packaging • units/pallet • storage area • safety stock", .65, 1.55, 3.75, 2.0, BLUE)
card(s, "SHIPMENT SCHEDULE", "shipment_id • product_id • warehouse_id • scheduled date • planned/actual pallets • type • status", 4.78, 1.55, 3.75, 2.0, CYAN)
card(s, "INVENTORY TRANSACTIONS", "transaction_id • shipment_id • product_id • warehouse_id • date • type • signed pallet quantity", 8.9, 1.55, 3.75, 2.0, ORANGE)
text(s, "JOIN KEYS", .75, 4.1, 1.4, .3, 11, CYAN, True)
text(s, "product_id  +  warehouse_id  +  shipment_id  +  date", .75, 4.47, 8, .55, 21, NAVY, True)
bullets(s, ["Synthetic records protect customers, employees, products, and shipment details.", "The same schema supports a 25,000-row demo and a 100-million-record cloud test."], .75, 5.15, 11.5, 1.2, 17)

# 4
s = new_slide("Projected inventory is calculated in event order", "Core calculation", "For each product and warehouse, I calculate current inventory from completed transactions, then apply future inbound and outbound events by date. Inbound is processed before outbound on the same date.")
text(s, "Projected inventory", .8, 1.55, 3.0, .5, 19, NAVY, True)
text(s, "= current inventory + cumulative inbound − cumulative outbound", 3.6, 1.55, 8.8, .5, 20, BLUE, True)
card(s, "TODAY", "20 pallets available", .8, 2.55, 2.55, 1.45, BLUE)
text(s, "→", 3.47, 2.93, .45, .35, 24, MUTED, True, PP_ALIGN.CENTER)
card(s, "OUTBOUND: 15", "5 pallets remain", 4.0, 2.55, 2.55, 1.45, GREEN)
text(s, "→", 6.67, 2.93, .45, .35, 24, MUTED, True, PP_ALIGN.CENTER)
card(s, "INBOUND: 10", "15 pallets available", 7.2, 2.55, 2.55, 1.45, CYAN)
text(s, "→", 9.87, 2.93, .45, .35, 24, MUTED, True, PP_ALIGN.CENTER)
card(s, "OUTBOUND: 18", "Short by 3 pallets", 10.4, 2.55, 2.15, 1.45, RED)
text(s, "Missing pallets = max(order quantity − inventory available before order, 0)", .85, 4.65, 11.6, .7, 19, NAVY, True, PP_ALIGN.CENTER)

# 5
s = new_slide("Every outbound order receives an actionable status", "Decision logic", "The status is designed for supervisors and planners. Ready and At Risk distinguish healthy orders from orders that leave little safety stock. Waiting for Inbound separates timing issues from unresolved shortages.")
card(s, "READY", "Enough inventory exists and the remaining balance stays at or above safety stock.", .7, 1.55, 5.8, 1.55, GREEN)
card(s, "AT RISK", "The order can be filled, but remaining inventory falls below the product safety-stock threshold.", 6.8, 1.55, 5.8, 1.55, ORANGE)
card(s, "WAITING FOR INBOUND", "The order is short on its scheduled date, but a later planned inbound is expected to cover it.", .7, 3.55, 5.8, 1.55, BLUE)
card(s, "SHORT", "The order lacks inventory and no known inbound covers the shortage within the planning horizon.", 6.8, 3.55, 5.8, 1.55, RED)
text(s, "The earliest-available date is an expectation—not a guarantee—because inbound shipments can be delayed.", .85, 5.75, 11.6, .55, 15, MUTED, False, PP_ALIGN.CENTER)

# 6
s = new_slide("Spark converts raw records into trusted analytical data", "Processing pipeline", "The Spark job reads explicit schemas, applies data-quality rules, calculates current inventory, orders future events with window functions, assigns readiness status, and writes partitioned Parquet.")
steps = [("1", "INGEST", "Explicit CSV schemas"), ("2", "VALIDATE", "Keys, dates, quantities"), ("3", "DEDUPLICATE", "Transaction & shipment IDs"), ("4", "JOIN", "Products, balances, schedule"), ("5", "WINDOW", "Cumulative event inventory"), ("6", "PUBLISH", "Partitioned Parquet")]
for i, (num, name, detail) in enumerate(steps):
    x = .55 + i*2.1
    card(s, f"{num}  {name}", detail, x, 2.0, 1.85, 1.55, [BLUE,CYAN,ORANGE,GREEN,BLUE,CYAN][i])
    if i < 5: text(s, "→", x+1.85, 2.55, .25, .25, 17, MUTED, True, PP_ALIGN.CENTER)
bullets(s, ["Invalid and unrealistic rows are written to quarantine rather than silently discarded.", "Adaptive query execution and partitioned output support scale without changing business logic.", "Parquet with Snappy compression reduces scan size for Spark and Athena."], .75, 4.25, 11.8, 1.75, 17)

# 7
s = new_slide("Local development and AWS use the same pipeline", "Platform architecture", "I first test the logic locally with Docker and a smaller dataset. The cloud version stores raw and curated zones in S3, runs Spark on EMR Serverless, and exposes results through Athena.")
arch = [("Synthetic data", "CSV partitions"), ("Amazon S3", "Raw zone"), ("EMR Serverless", "Apache Spark"), ("Amazon S3", "Curated Parquet"), ("Amazon Athena", "SQL analysis"), ("Dashboard", "Operational KPIs")]
for i, (name, detail) in enumerate(arch):
    x = .5 + i*2.12
    card(s, name, detail, x, 2.0, 1.82, 1.6, [BLUE,CYAN,ORANGE,CYAN,BLUE,GREEN][i])
    if i < 5: text(s, "→", x+1.84, 2.58, .28, .25, 17, MUTED, True, PP_ALIGN.CENTER)
text(s, "Terraform provisions the S3 bucket, Glue database, Athena workgroup, IAM role, and EMR Serverless application.", .75, 4.35, 11.8, .65, 18, NAVY, True, PP_ALIGN.CENTER)
text(s, "Local: Docker + Spark  |  Cloud: S3 + EMR Serverless + Glue/Athena  |  Format: Snappy Parquet", .75, 5.45, 11.8, .5, 15, MUTED, False, PP_ALIGN.CENTER)

# 8
s = new_slide("The reproducible sample demonstrates the full workflow", "Synthetic data", f"The included sample uses a fixed random seed and an as-of date of August 1, 2026. It represents three warehouses and realistic inbound, outbound, transfer, put-away, adjustment, and delay patterns.")
metrics = [(f"{len(transactions):,}", "inventory transactions"), (f"{len(shipments):,}", "scheduled shipments"), (f"{len(products):,}", "products"), ("3", "warehouses")]
for i, (value, label) in enumerate(metrics):
    x = .7 + i*3.05
    card(s, value, label, x, 1.7, 2.65, 1.4, [BLUE,CYAN,ORANGE,GREEN][i])
bullets(s, ["Historical completed shipments create the current inventory balance.", "Future scheduled and delayed shipments create the planning horizon.", "Opening balances and adjustments make shortage scenarios realistic without exposing company data.", "A command-line row-count parameter supports the 100-million-record requirement."], .8, 3.65, 11.6, 2.25, 18)

# 9
s = new_slide("The sample identifies 173 orders requiring attention", "Demonstration results", f"Of {len(readiness):,} upcoming outbound orders, {statuses['READY']:,} are ready. The remaining {len(readiness)-statuses['READY']:,} are at risk, waiting for inbound, or short. Total modeled shortage exposure is {missing_total:,} pallets.")
s.shapes.add_picture(str(ART / "readiness_status.png"), Inches(.7), Inches(1.45), width=Inches(7.7))
card(s, "ORDERS EVALUATED", f"{len(readiness):,}", 9.0, 1.65, 3.2, 1.15, BLUE)
card(s, "NEED ATTENTION", f"{len(readiness)-statuses['READY']:,}", 9.0, 3.0, 3.2, 1.15, ORANGE)
card(s, "MISSING PALLETS", f"{missing_total:,}", 9.0, 4.35, 3.2, 1.15, RED)
text(s, "Synthetic proof-of-concept results", 9.0, 5.75, 3.2, .3, 11, MUTED, False, PP_ALIGN.CENTER)

# 10
s = new_slide("Athena answers operational questions without rescanning CSV", "Analytics", "The curated tables support daily exception reporting and broader product analysis. These example results are produced from the same synthetic schedule included with the project.")
s.shapes.add_picture(str(ART / "product_volume.png"), Inches(.6), Inches(1.45), width=Inches(7.6))
bullets(s, ["Which orders are short or waiting for inbound?", "How many pallets are missing by warehouse?", "What is the first expected availability date?", "Which products drive the most inbound and outbound volume?", "Which products experience the most adjustments and delays?"], 8.6, 1.65, 3.9, 3.7, 16)
text(s, "Partition pruning + columnar storage = lower query cost and faster dashboards", 8.65, 5.55, 3.75, .65, 14, BLUE, True, PP_ALIGN.CENTER)

# 11
s = new_slide("The design scales to 100 million records", "Big-data strategy", "The large test uses the same schemas and logic, but generates partitioned files and distributes processing across Spark executors. The goal is to demonstrate volume, distributed computation, and efficient query storage—not merely create one oversized CSV.")
card(s, "GENERATE IN PARTS", "Create deterministic partitions so generation can restart and parallelize.", .7, 1.55, 3.75, 1.55, BLUE)
card(s, "DISTRIBUTE BY KEY", "Repartition calculations by warehouse and product; avoid collecting data to the driver.", 4.78, 1.55, 3.75, 1.55, CYAN)
card(s, "STORE COLUMNAR", "Write Snappy Parquet partitioned by warehouse and event month.", 8.87, 1.55, 3.75, 1.55, ORANGE)
bullets(s, ["Use Spark window functions for chronological balances.", "Enable adaptive query execution to improve shuffle plans.", "Inspect row counts, rejected records, partition sizes, and data skew.", "Use Athena workgroup controls and S3 lifecycle policies to manage cost."], .8, 3.75, 11.5, 2.15, 18)

# 12
s = new_slide("The project turns warehouse events into early action", "Conclusion", "The expected result is a practical analytics system that finds inventory constraints before orders are loaded. It demonstrates data modeling, quality controls, distributed processing, cloud design, SQL analytics, and business interpretation.")
card(s, "OPERATIONAL VALUE", "Earlier shortage visibility • quantified pallet gaps • clearer inbound dependency", .7, 1.55, 3.75, 1.7, GREEN)
card(s, "TECHNICAL VALUE", "Spark windows • Parquet • Athena SQL • Terraform • reproducibility", 4.78, 1.55, 3.75, 1.7, BLUE)
card(s, "NEXT STEPS", "Run the 100M test • capture performance metrics • connect a dashboard", 8.87, 1.55, 3.75, 1.7, ORANGE)
text(s, "Warehouse Inventory and Outbound Order Readiness", .75, 4.15, 11.8, .6, 25, NAVY, True, PP_ALIGN.CENTER)
text(s, "Identify the shortage before the loading appointment.", .75, 5.05, 11.8, .5, 18, CYAN, True, PP_ALIGN.CENTER)
text(s, "Thank you — Questions?", .75, 6.0, 11.8, .5, 20, INK, True, PP_ALIGN.CENTER)

output = ART / "CS675_Warehouse_Inventory_Order_Readiness_Shuhao_Lin.pptx"
prs.save(output)

narration = ART / "CS675_Presentation_Speaker_Notes.md"
with narration.open("w", encoding="utf-8") as handle:
    handle.write("# CS-675 Presentation Speaker Notes\n\n")
    for number, slide in enumerate(prs.slides, 1):
        title = next((shape.text.strip().replace("\n", " ") for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()), f"Slide {number}")
        note = slide.notes_slide.notes_text_frame.text.strip()
        handle.write(f"## Slide {number}: {title}\n\n{note}\n\n")

print(f"Created {output} with {len(prs.slides)} slides")
